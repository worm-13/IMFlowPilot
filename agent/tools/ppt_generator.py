import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)


COLOR_SCHEMES = {
    "商务蓝": {
        "primary": RGBColor(0, 112, 192),
        "secondary": RGBColor(68, 114, 196),
        "background": RGBColor(255, 255, 255),
        "text": RGBColor(0, 0, 0),
        "accent": RGBColor(79, 129, 189),
        "light_bg": RGBColor(217, 226, 243),
    },
    "商务红": {
        "primary": RGBColor(192, 0, 0),
        "secondary": RGBColor(217, 80, 68),
        "background": RGBColor(255, 255, 255),
        "text": RGBColor(0, 0, 0),
        "accent": RGBColor(255, 192, 0),
        "light_bg": RGBColor(255, 235, 234),
    },
    "商务绿": {
        "primary": RGBColor(0, 128, 0),
        "secondary": RGBColor(112, 173, 71),
        "background": RGBColor(255, 255, 255),
        "text": RGBColor(0, 0, 0),
        "accent": RGBColor(146, 208, 80),
        "light_bg": RGBColor(226, 239, 218),
    },
    "深蓝科技": {
        "primary": RGBColor(0, 32, 96),
        "secondary": RGBColor(0, 112, 192),
        "background": RGBColor(30, 30, 30),
        "text": RGBColor(255, 255, 255),
        "accent": RGBColor(0, 112, 192),
        "light_bg": RGBColor(50, 50, 80),
    },
    "商务紫": {
        "primary": RGBColor(112, 48, 160),
        "secondary": RGBColor(146, 73, 168),
        "background": RGBColor(255, 255, 255),
        "text": RGBColor(0, 0, 0),
        "accent": RGBColor(177, 145, 200),
        "light_bg": RGBColor(238, 227, 244),
    },
}


class PPTGenerator:

    @staticmethod
    def generate_ppt(topic: str, slides_content: list[dict], theme: str = "商务蓝") -> str:
        if theme not in COLOR_SCHEMES:
            theme = "商务蓝"

        colors = COLOR_SCHEMES[theme]
        prs = Presentation()

        for idx, slide_data in enumerate(slides_content):
            is_cover_or_ending = (idx == 0 or idx == len(slides_content) - 1)
            is_toc_or_section = ("目录" in slide_data["title"] or "章节" in slide_data["title"])

            if is_cover_or_ending:
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = slide_data["title"]
                title.text_frame.paragraphs[0].font.color.rgb = colors["primary"]
                title.text_frame.paragraphs[0].font.size = Pt(44)
                title.text_frame.paragraphs[0].font.bold = True

                if len(slide.placeholders) > 1 and slide.placeholders[1].has_text_frame:
                    subtitle = slide.placeholders[1]
                    subtitle.text = slide_data["content"]
                    subtitle.text_frame.paragraphs[0].font.color.rgb = colors["secondary"]
                    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
                    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            elif is_toc_or_section:
                slide_layout = prs.slide_layouts[2]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = slide_data["title"]
                title.text_frame.paragraphs[0].font.color.rgb = colors["primary"]
                title.text_frame.paragraphs[0].font.size = Pt(40)
                title.text_frame.paragraphs[0].font.bold = True

                if len(slide.placeholders) > 1 and slide.placeholders[1].has_text_frame:
                    content = slide.placeholders[1]
                    content.text = slide_data["content"]
                    content.text_frame.paragraphs[0].font.color.rgb = colors["text"]

            else:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]

                title.text = slide_data["title"]
                title.text_frame.paragraphs[0].font.color.rgb = colors["primary"]
                title.text_frame.paragraphs[0].font.size = Pt(36)
                title.text_frame.paragraphs[0].font.bold = True

                content.text = slide_data["content"]
                content.text_frame.paragraphs[0].font.color.rgb = colors["text"]
                content.text_frame.paragraphs[0].font.size = Pt(18)

        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        file_name = f"PPT_{timestamp}.pptx"
        file_path = os.path.join(OUTPUT_DIR, file_name)
        prs.save(file_path)
        return file_path